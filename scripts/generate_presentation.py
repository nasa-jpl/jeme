#!/usr/bin/env python3
"""Generate JEME/JEOE Technical Seminar slide deck."""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1E, 0x29, 0x3B)       # slate-800
GRAY700 = RGBColor(0x37, 0x41, 0x51)
GRAY500 = RGBColor(0x6B, 0x72, 0x80)
GRAY400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY200 = RGBColor(0xE5, 0xE7, 0xEB)
BLUE600 = RGBColor(0x25, 0x63, 0xEB)
BLUE500 = RGBColor(0x3B, 0x82, 0xF6)
BLUE100 = RGBColor(0xDB, 0xEA, 0xFE)
INDIGO600 = RGBColor(0x4F, 0x46, 0xE5)
INDIGO100 = RGBColor(0xE0, 0xE7, 0xFF)
GREEN600 = RGBColor(0x05, 0x9F, 0x46)
GREEN500 = RGBColor(0x10, 0xB9, 0x81)
GREEN100 = RGBColor(0xD1, 0xFA, 0xE5)
AMBER500 = RGBColor(0xF5, 0x9E, 0x0B)
AMBER100 = RGBColor(0xFE, 0xF3, 0xC7)
RED500 = RGBColor(0xEF, 0x44, 0x44)
PURPLE500 = RGBColor(0x8B, 0x5C, 0xF6)
CYAN500 = RGBColor(0x06, 0xB6, 0xD4)

MODEL_COLORS = {
    'RAPID': RGBColor(0x3B, 0x82, 0xF6),
    'CARDAMOM': RGBColor(0x10, 0xB9, 0x81),
    'CMS-Flux': RGBColor(0xF5, 0x9E, 0x0B),
    'ECCO': RGBColor(0x8B, 0x5C, 0xF6),
    'ISSM': RGBColor(0x06, 0xB6, 0xD4),
    'MOMO-CHEM': RGBColor(0xEF, 0x44, 0x44),
    'LES': RGBColor(0x2E, 0x8B, 0x57),
    'EDMF': RGBColor(0xFF, 0x63, 0x47),
}

SCORE_COLORS = {
    0: RGBColor(0xD1, 0xD5, 0xDB),
    1: RGBColor(0xFD, 0xE6, 0x8A),
    2: RGBColor(0x93, 0xC5, 0xFD),
    3: RGBColor(0x6E, 0xE7, 0xB7),
}

SCORE_LABELS = {0: 'None', 1: 'Low', 2: 'Medium', 3: 'High'}

# ── Helpers ─────────────────────────────────────────────────────────────────
def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Add a colored background rectangle."""
    w = width or Inches(13.33)
    h = height or Inches(7.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=GRAY700, spacing=Pt(6), bold_prefix=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        if bold_prefix and ':' in item:
            prefix, rest = item.split(':', 1)
            run1 = p.add_run()
            run1.text = prefix + ':'
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = '• ' + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
    return txBox


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_number_footer(slide, num, total):
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
                f'{num}/{total}', font_size=10, color=GRAY400, alignment=PP_ALIGN.RIGHT)


def section_header_slide(prs, title, subtitle, slide_num, total, accent_color=BLUE600):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, RGBColor(0xF8, 0xFA, 0xFC))
    add_accent_bar(slide, Inches(0), Inches(2.8), Inches(13.33), Inches(0.06), accent_color)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11), Inches(1.2),
                title, font_size=36, bold=True, color=DARK)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                subtitle, font_size=18, color=GRAY500)
    slide_number_footer(slide, slide_num, total)
    return slide


def colored_cell_table(slide, left, top, rows_data, col_widths, header_row=True):
    """Draw a manual table using shapes for better color control."""
    row_h = Inches(0.38)
    header_h = Inches(0.42)
    for r, row in enumerate(rows_data):
        x = left
        h = header_h if (r == 0 and header_row) else row_h
        y_offset = top + (header_h if header_row else 0) + row_h * (r - (1 if header_row else 0)) if r > 0 else top
        if r == 0:
            y_offset = top
        else:
            y_offset = top + header_h + row_h * (r - 1) if header_row else top + row_h * r

        for c, cell in enumerate(row):
            w = col_widths[c]
            text = cell.get('text', '')
            bg = cell.get('bg', None)
            fg = cell.get('fg', DARK)
            bold = cell.get('bold', False)
            fs = cell.get('fs', 11)
            align = cell.get('align', PP_ALIGN.CENTER)

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_offset, w, h)
            if bg:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg
            else:
                shape.fill.background()
            shape.line.color.rgb = GRAY200
            shape.line.width = Pt(0.5)

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]
            p.text = str(text)
            p.font.size = Pt(fs)
            p.font.bold = bold
            p.font.color.rgb = fg
            p.font.name = 'Calibri'
            p.alignment = align

            x += w


# ── Load MCL data ──────────────────────────────────────────────────────────
script_dir = os.path.dirname(os.path.abspath(__file__))
project_dir = os.path.dirname(script_dir)
mcl_path = os.path.join(project_dir, 'src', 'data', 'mcl_scores.json')
with open(mcl_path) as f:
    mcl_data = json.load(f)

TIER1_DIMS = [
    ('process_representation', 'Process Representation'),
    ('spatial_resolution', 'Spatial Resolution'),
    ('temporal_resolution', 'Temporal Resolution'),
    ('process_coupling', 'Process Coupling'),
    ('predictive_skill', 'Predictive Skill'),
    ('computational_performance', 'Computational Performance'),
    ('observational_constraint', 'Observational Constraint'),
    ('retrospective_analysis', 'Retrospective Analysis'),
    ('uq_attribution', 'Uncertainty Quantification & Attribution'),
    ('vv_framework', 'Verification & Validation Framework'),
    ('ml_ai_integration', 'Machine Learning / Artificial Intelligence Integration'),
    ('mission_support', 'Mission Support'),
    ('interoperability', 'Interoperability'),
    ('stakeholder_adoption', 'Stakeholder Adoption'),
]

MODELS = ['ECCO', 'ISSM', 'CMS-Flux', 'RAPID', 'CARDAMOM', 'MOMO-CHEM', 'LES', 'EDMF']

TOTAL_SLIDES = 15

# ── Build Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, RGBColor(0x0F, 0x17, 0x2A))  # dark navy

# Accent lines
add_accent_bar(sl, Inches(1), Inches(2.0), Inches(4), Inches(0.05), BLUE500)
add_accent_bar(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.03), RGBColor(0x33, 0x44, 0x66))

# Title
add_textbox(sl, Inches(1), Inches(2.3), Inches(11), Inches(1.8),
            "JPL's Earth Modeling & Observation Dashboard",
            font_size=38, bold=True, color=WHITE)
add_textbox(sl, Inches(1), Inches(4.0), Inches(11), Inches(1.0),
            "Tracking Scientific Impact, Cross-Model Connections, and Capability Assessment",
            font_size=22, color=RGBColor(0x93, 0xC5, 0xFD))
add_textbox(sl, Inches(1), Inches(5.8), Inches(5), Inches(0.5),
            "Technical Seminar  •  March 2026", font_size=14, color=GRAY400)
add_textbox(sl, Inches(1), Inches(6.3), Inches(5), Inches(0.5),
            "JPL's Earth Modeling Enterprise  •  JPL Earth Observation Enterprise  •  Model Capability Levels", font_size=12, color=RGBColor(0x64, 0x74, 0x8B))
slide_number_footer(sl, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(0.12), Inches(7.5), BLUE600)

add_textbox(sl, Inches(0.7), Inches(0.4), Inches(10), Inches(0.7),
            "Agenda", font_size=32, bold=True, color=DARK)

toc_items = [
    ("1", "Motivation & Objectives", "Why quantify model impact and maturity?", "1.5 min", BLUE600),
    ("2", "JEME Overview", "8 Earth system models, domains, coupling", "1.5 min", BLUE500),
    ("3", "JEOE Overview", "NASA missions (GRACE, SWOT), relationship to models", "1 min", PURPLE500),
    ("4-5", "Citation Analytics Pipeline", "Data collection, Large Language Model classification, verification", "2.5 min", GREEN500),
    ("6", "Data Quality & Uncertainty", "Multi-agent verification, 3-phase uncertainty quantification pipeline", "1.5 min", AMBER500),
    ("7", "Cross-Model Network Analysis", "Bridge papers, connection matrix, co-authorship", "1.5 min", CYAN500),
    ("8", "Earth System Interconnections", "Sphere classification, inter-sphere research links", "2 min", GREEN600),
    ("9", "Model Capability Level Framework", "Heritage (Technology Readiness Levels → Predictive Capability Maturity Model → Model Capability Levels), Tier 1 + Tier 2", "2 min", INDIGO600),
    ("10-11", "Model Capability Level Results", "Radar chart, heatmap, model deep-dives", "2.5 min", INDIGO600),
    ("12", "Gap Analysis", "Common development needs across JEME", "1.5 min", RED500),
    ("13", "Dashboard Architecture", "React, data flow, interactive features", "1 min", GRAY500),
    ("14-15", "Path Forward & Summary", "Large Language Model-enhanced scoring, next steps, takeaways", "1.5 min", DARK),
]

y = Inches(1.3)
for num, title, desc, time, color in toc_items:
    # Number badge
    badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(0.55), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    add_textbox(sl, Inches(1.45), y - Pt(1), Inches(5.5), Inches(0.3),
                title, font_size=13, bold=True, color=DARK)
    add_textbox(sl, Inches(1.45), y + Pt(14), Inches(5.5), Inches(0.25),
                desc, font_size=10, color=GRAY500)
    add_textbox(sl, Inches(7.2), y + Pt(4), Inches(1), Inches(0.3),
                time, font_size=10, color=GRAY400, alignment=PP_ALIGN.RIGHT)

    y += Inches(0.48)

slide_number_footer(sl, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVATION
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), BLUE600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Motivation & Objectives", font_size=30, bold=True, color=DARK)

# Three columns
col_data = [
    ("Scientific Impact", BLUE600,
     ["How widely are JPL models cited?",
      "Which research domains adopt them?",
      "What is the geographic reach?",
      "Track growth and influence over time"]),
    ("Cross-Model Connections", GREEN600,
     ["Which models share citing papers?",
      "How do Earth system spheres interconnect?",
      "Where are collaboration opportunities?",
      "Identify bridge papers linking domains"]),
    ("Capability Assessment", INDIGO600,
     ["Are models fit for intended applications?",
      "Where are common development gaps?",
      "How to prioritize limited resources?",
      "Systematic, evidence-based evaluation"]),
]

for i, (title, color, items) in enumerate(col_data):
    x = Inches(0.8 + i * 4.0)
    add_accent_bar(sl, x, Inches(1.3), Inches(3.5), Inches(0.05), color)
    add_textbox(sl, x, Inches(1.55), Inches(3.5), Inches(0.5),
                title, font_size=18, bold=True, color=color)
    add_bullet_list(sl, x, Inches(2.2), Inches(3.5), Inches(3.5),
                    items, font_size=13, color=GRAY700)

# Bottom context
add_textbox(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(1.2),
            "Context: NASA's Earth Science to Action Strategy requires objective assessment of observation-informed\n"
            "models for advancing science, developing decision-support tools, and formulating new satellite missions.",
            font_size=12, color=GRAY500)

slide_number_footer(sl, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — JEME OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), BLUE600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "JEME: JPL's Earth Modeling Enterprise", font_size=30, bold=True, color=DARK)
add_textbox(sl, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
            "8 observation-informed modeling frameworks spanning the full Earth system", font_size=14, color=GRAY500)

model_info = [
    ('ECCO', 'Oceanography', 'Four-Dimensional Variational adjoint', 'Massachusetts Institute of Technology, University of Texas Austin, Scripps Institution of Oceanography'),
    ('ISSM', 'Cryosphere / Sea Level', 'Four-Dimensional Variational adjoint', 'University of California, Irvine'),
    ('CMS-Flux', 'Carbon Flux Monitoring', 'Four-Dimensional Variational adjoint', 'NASA Goddard Space Flight Center, Ames'),
    ('RAPID', 'Hydrology', 'Kalman Filter', 'University of Texas Austin, National Oceanic and Atmospheric Administration National Weather Service'),
    ('CARDAMOM', 'Terrestrial Ecosystems', 'Bayesian Markov Chain Monte Carlo', 'Stanford, Edinburgh'),
    ('MOMO-CHEM', 'Atmospheric Chemistry', 'Ensemble Kalman Filter', 'Japan Agency for Marine-Earth Science and Technology'),
    ('LES', 'Cloud / Turbulence', 'Process model', 'University of Connecticut'),
    ('EDMF', 'Atmospheric Physics', 'Bayesian Inference', 'Geophysical Fluid Dynamics Laboratory, National Center for Atmospheric Research'),
]

headers = [
    {'text': 'Model', 'bg': RGBColor(0x1E, 0x40, 0xAF), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Domain', 'bg': RGBColor(0x1E, 0x40, 0xAF), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Data Assimilation Method', 'bg': RGBColor(0x1E, 0x40, 0xAF), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Key Partners', 'bg': RGBColor(0x1E, 0x40, 0xAF), 'fg': WHITE, 'bold': True, 'fs': 11},
]

rows = [headers]
for name, domain, da, partners in model_info:
    mc = MODEL_COLORS.get(name, GRAY500)
    rows.append([
        {'text': name, 'fg': mc, 'bold': True, 'fs': 11, 'align': PP_ALIGN.LEFT},
        {'text': domain, 'fg': GRAY700, 'fs': 11, 'align': PP_ALIGN.LEFT},
        {'text': da, 'fg': GRAY700, 'fs': 11, 'align': PP_ALIGN.LEFT},
        {'text': partners, 'fg': GRAY500, 'fs': 9, 'align': PP_ALIGN.LEFT},
    ])

col_w = [Inches(1.6), Inches(2.2), Inches(3.2), Inches(5.0)]
colored_cell_table(sl, Inches(0.5), Inches(1.6), rows, col_w)

add_textbox(sl, Inches(0.8), Inches(5.8), Inches(11), Inches(1.0),
            "All models are mission-essential JPL assets for: (1) addressing Decadal Survey questions,\n"
            "(2) developing decision-support applications, and (3) new satellite mission formulation.",
            font_size=12, color=GRAY500)

slide_number_footer(sl, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — JEOE OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), PURPLE500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "JEOE: JPL Earth Observation Enterprise", font_size=30, bold=True, color=DARK)

# GRACE card
for i, (name, desc, domain, highlights) in enumerate([
    ("GRACE / GRACE Follow-On", "Gravity Recovery and Climate Experiment",
     "Geophysics & Geodesy",
     ["Tracks Earth's gravity field changes", "Water storage, ice mass balance, sea level",
      "Provides observational constraint for ECCO, ISSM, RAPID"]),
    ("SWOT", "Surface Water and Ocean Topography",
     "Ocean & Inland Water",
     ["Ka-band radar interferometry", "Unprecedented water surface elevation measurements",
      "Key future data source for RAPID, ECCO"]),
]):
    x = Inches(0.8 + i * 6.0)
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(5.2), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    box.line.color.rgb = GRAY200

    add_accent_bar(sl, x, Inches(1.5), Inches(5.2), Inches(0.06), PURPLE500)
    add_textbox(sl, x + Inches(0.3), Inches(1.8), Inches(4.6), Inches(0.4),
                name, font_size=20, bold=True, color=DARK)
    add_textbox(sl, x + Inches(0.3), Inches(2.2), Inches(4.6), Inches(0.3),
                desc, font_size=11, color=GRAY500)
    add_textbox(sl, x + Inches(0.3), Inches(2.6), Inches(4.6), Inches(0.3),
                domain, font_size=12, bold=True, color=PURPLE500)
    add_bullet_list(sl, x + Inches(0.3), Inches(3.1), Inches(4.6), Inches(1.8),
                    highlights, font_size=12, color=GRAY700)

add_textbox(sl, Inches(0.8), Inches(5.5), Inches(11), Inches(1.0),
            "JEOE and JEME share a unified dashboard with context-aware navigation, cross-referencing\n"
            "mission observations with the models that assimilate their data.",
            font_size=12, color=GRAY500)

slide_number_footer(sl, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CITATION ANALYTICS PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), GREEN500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Citation Analytics Pipeline", font_size=30, bold=True, color=DARK)

# Pipeline stages as connected boxes
stages = [
    ("1. Seed Papers", "Identify core\nalgorithm &\napplication papers", BLUE600),
    ("2. Citation\nCollection", "Semantic Scholar\nCrossref APIs\nbatch retrieval", GREEN600),
    ("3. Large Language\nModel Classification", "Gemma3/Ollama:\nengagement level,\ndomain, geography", PURPLE500),
    ("4. Multi-Agent\nVerification", "Crossref +\nSemantic Scholar +\nkeyword classifier +\ndeduplication agent", AMBER500),
    ("5. Dashboard\nVisualization", "React + Recharts\ninteractive charts\nper-model views", RED500),
]

for i, (title, desc, color) in enumerate(stages):
    x = Inches(0.4 + i * 2.5)
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.1), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    add_textbox(sl, x + Inches(0.15), Inches(1.6), Inches(1.8), Inches(0.8),
                title, font_size=12, bold=True, color=WHITE)
    add_textbox(sl, x + Inches(0.15), Inches(2.4), Inches(1.8), Inches(1.2),
                desc, font_size=10, color=RGBColor(0xE0, 0xE0, 0xE0))

    # Arrow connector
    if i < len(stages) - 1:
        add_textbox(sl, x + Inches(2.1), Inches(2.3), Inches(0.4), Inches(0.4),
                    "→", font_size=20, bold=True, color=GRAY400, alignment=PP_ALIGN.CENTER)

# Key metrics
add_textbox(sl, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
            "Key Metrics per Model", font_size=16, bold=True, color=DARK)

metrics = [
    "Total papers & citations with growth trends",
    "Engagement levels: Level 1 (background) → Level 4 (foundational)",
    "Research domain distribution & geographic reach",
    "Cross-model bridge papers & co-authorship networks"
]
add_bullet_list(sl, Inches(0.8), Inches(4.8), Inches(11), Inches(2.0),
                metrics, font_size=13, color=GRAY700)

slide_number_footer(sl, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATA QUALITY & UQ
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), AMBER500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Data Quality & Uncertainty Quantification", font_size=30, bold=True, color=DARK)

# Two columns
# Left: verification
add_textbox(sl, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "Multi-Agent Verification", font_size=18, bold=True, color=AMBER500)
verify_items = [
    "Team Paper Categorizer: relevance tier classification",
    "Crossref Agent: Digital Object Identifier validation & venue metadata",
    "Semantic Scholar Agent: title recovery, batch enrichment",
    "Keyword Classifier: domain-specific relevance scoring",
    "Deduplication Agent: Digital Object Identifier-first, title-fallback detection",
    "ECCO: removed ~3,900 off-topic entries; enriched 7,600+ venues",
    "ISSM: repaired 1,904 broken team paper titles"
]
add_bullet_list(sl, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.0),
                verify_items, font_size=12, color=GRAY700)

# Right: UQ pipeline
add_textbox(sl, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
            "3-Phase Uncertainty Pipeline", font_size=18, bold=True, color=INDIGO600)
uq_items = [
    "Phase 1: Deterministic scoring — metadata completeness, keyword agreement, composite confidence",
    "Phase 2: Multi-temperature Large Language Model sampling — Gemini 3× at temperatures [0.1, 0.5, 1.0], stochastic variance",
    "Phase 3: Skeptic agent — challenges high-risk entries, override flags for disagreement ≤ 2/5",
    "Result: calibrated confidence scores on every citation entry",
    "False removal rate: ~8-12%; false retention: ~3-5%"
]
add_bullet_list(sl, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.0),
                uq_items, font_size=12, color=GRAY700)

slide_number_footer(sl, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NETWORK ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), CYAN500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Cross-Model Network Analysis", font_size=30, bold=True, color=DARK)

add_textbox(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
            "What We Compute", font_size=18, bold=True, color=CYAN500)
network_items = [
    "Bridge papers: shared citations across multiple models",
    "Connection matrix: pairwise model overlap counts",
    "Cross-model co-authorship: researchers spanning model teams",
    "Domain overlap: shared research domain classifications",
    "Network graph: force-directed visualization of connectivity"
]
add_bullet_list(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0),
                network_items, font_size=13, color=GRAY700)

add_textbox(sl, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
            "Key Insights", font_size=18, bold=True, color=CYAN500)
insight_items = [
    "Strongest connections: Estimating the Circulation and Climate of the Ocean ↔ Ice Sheet System Model (ocean-ice coupling), Carbon Monitoring System Flux ↔ CARDAMOM (carbon cycle)",
    "Bridge papers reveal interdisciplinary research opportunities",
    "Identifies potential collaboration gaps between modeling teams"
]
add_bullet_list(sl, Inches(0.8), Inches(5.1), Inches(11), Inches(2.0),
                insight_items, font_size=13, color=GRAY700)

# Visual placeholder for network graph
box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.2), Inches(5.5), Inches(3.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
box.line.color.rgb = GRAY200
add_textbox(sl, Inches(7.5), Inches(2.5), Inches(4.5), Inches(0.8),
            "[Network Graph Visualization]\nInteractive force-directed graph from dashboard",
            font_size=12, color=GRAY400, alignment=PP_ALIGN.CENTER)

slide_number_footer(sl, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EARTH SYSTEM INTERCONNECTIONS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), GREEN600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Earth System Interconnections", font_size=30, bold=True, color=DARK)
add_textbox(sl, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
            "Classifying all JEME publications into Earth's five interconnected spheres", font_size=14, color=GRAY500)

sphere_data = [
    ("Atmosphere", RGBColor(0x60, 0xA5, 0xFA), "Clouds, radiation, chemistry, transport"),
    ("Hydrosphere", RGBColor(0x38, 0xBD, 0xF8), "Rivers, oceans, water cycle, precipitation"),
    ("Cryosphere", RGBColor(0x7D, 0xD3, 0xFC), "Ice sheets, glaciers, sea ice, permafrost"),
    ("Biosphere", RGBColor(0x4A, 0xDE, 0x80), "Vegetation, carbon cycle, ecosystems"),
    ("Geosphere", RGBColor(0xFB, 0xBF, 0x24), "Solid Earth, tectonics, gravity, topography"),
]

for i, (name, color, desc) in enumerate(sphere_data):
    x = Inches(0.4 + i * 2.5)
    box = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(1.8), Inches(1.2), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = name[:4]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(18)

    add_textbox(sl, x, Inches(3.1), Inches(2.3), Inches(0.3),
                name, font_size=12, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, Inches(3.4), Inches(2.3), Inches(0.8),
                desc, font_size=10, color=GRAY500, alignment=PP_ALIGN.CENTER)

# Features
add_textbox(sl, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
            "Analysis Features", font_size=16, bold=True, color=GREEN600)
earth_items = [
    "Keyword-based paper classification into spheres (multi-sphere papers counted for each)",
    "Inter-sphere connection matrix: papers bridging multiple Earth system components",
    "Interactive D3 force-directed graph with drag, zoom, hover highlighting, animated particles",
    "Model-to-sphere mapping: which JEME models contribute to which spheres",
    "Cross-sphere bridge papers table: highest-impact interdisciplinary research"
]
add_bullet_list(sl, Inches(0.8), Inches(5.1), Inches(11), Inches(2.0),
                earth_items, font_size=12, color=GRAY700)

slide_number_footer(sl, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MCL FRAMEWORK
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), INDIGO600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Model Capability Level (MCL) Framework", font_size=30, bold=True, color=DARK)

# Heritage arrow
heritage = [("NASA Technology\nReadiness Levels", "Technology\nreadiness\n(9 levels)", BLUE600),
            ("Department of Energy\nPredictive Capability\nMaturity Model", "Predictive\ncapability\n(4 levels)", AMBER500),
            ("JPL Model\nCapability Levels", "Earth system\nmodel assessment\n(0-3 scale)", INDIGO600)]

for i, (title, desc, color) in enumerate(heritage):
    x = Inches(0.8 + i * 3.8)
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.0), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    add_textbox(sl, x + Inches(0.2), Inches(1.4), Inches(2.6), Inches(0.4),
                title, font_size=16, bold=True, color=WHITE)
    add_textbox(sl, x + Inches(0.2), Inches(1.85), Inches(2.6), Inches(0.8),
                desc, font_size=11, color=RGBColor(0xE0, 0xE0, 0xE0))
    if i < 2:
        add_textbox(sl, x + Inches(3.0), Inches(1.7), Inches(0.8), Inches(0.5),
                    "→", font_size=24, bold=True, color=GRAY400, alignment=PP_ALIGN.CENTER)

# Tier 1 and Tier 2 summary
add_textbox(sl, Inches(0.8), Inches(3.2), Inches(5.5), Inches(0.5),
            "Tier 1: Core Capabilities (14 dimensions)", font_size=16, bold=True, color=INDIGO600)
t1_items = [
    "Process Representation, Spatial & Temporal Resolution",
    "Process Coupling, Predictive Skill, Computational Performance",
    "Observational Constraint, Retrospective Analysis",
    "Uncertainty Quantification & Attribution, Verification & Validation Framework, Machine Learning / Artificial Intelligence Integration",
    "Mission Support, Interoperability, Stakeholder Adoption"
]
add_bullet_list(sl, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.5),
                t1_items, font_size=12, color=GRAY700)

add_textbox(sl, Inches(7.0), Inches(3.2), Inches(5.5), Inches(0.5),
            "Tier 2: Application Domains (5 × 4 sub-dimensions)", font_size=16, bold=True, color=INDIGO600)
t2_items = [
    "Research: process fidelity, coupling, community impact, extremes",
    "Data Assimilation: observational coverage, error correction, operational use",
    "Verification, Validation, Uncertainty Quantification / Machine Learning: verification, artificial intelligence integration, emulator validation",
    "Mission Support: Observing System Simulation Experiments, data assimilation readiness, Decadal Survey alignment",
    "Decision Support: usability, regulatory relevance, uncertainty quantification communication"
]
add_bullet_list(sl, Inches(7.0), Inches(3.8), Inches(5.5), Inches(2.5),
                t2_items, font_size=12, color=GRAY700)

# Scoring scale
y = Inches(6.4)
add_textbox(sl, Inches(0.8), y, Inches(1.5), Inches(0.3),
            "Scoring Scale:", font_size=11, bold=True, color=DARK)
for level in range(4):
    x = Inches(2.5 + level * 2.5)
    badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = SCORE_COLORS[level]
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"Level {level} — {SCORE_LABELS[level]}"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

slide_number_footer(sl, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MCL RESULTS: HEATMAP
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), INDIGO600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "MCL Results: Tier 1 Capability Heatmap", font_size=30, bold=True, color=DARK)

# Build heatmap table
header = [{'text': 'Model', 'bg': RGBColor(0x31, 0x2E, 0x81), 'fg': WHITE, 'bold': True, 'fs': 9}]
for _, label in TIER1_DIMS:
    header.append({'text': label, 'bg': RGBColor(0x31, 0x2E, 0x81), 'fg': WHITE, 'bold': True, 'fs': 6})
header.append({'text': 'Average', 'bg': RGBColor(0x31, 0x2E, 0x81), 'fg': WHITE, 'bold': True, 'fs': 8})

hm_rows = [header]
for model in MODELS:
    row = [{'text': model, 'fg': MODEL_COLORS.get(model, DARK), 'bold': True, 'fs': 10, 'align': PP_ALIGN.LEFT}]
    scores = []
    for dim_id, _ in TIER1_DIMS:
        s = mcl_data.get(model, {}).get('tier1', {}).get(dim_id, {}).get('score', 0)
        scores.append(s)
        row.append({'text': str(s), 'bg': SCORE_COLORS.get(s, SCORE_COLORS[0]), 'fg': DARK, 'bold': False, 'fs': 10})
    avg = sum(scores) / len(scores) if scores else 0
    row.append({'text': f'{avg:.1f}', 'fg': DARK, 'bold': False, 'fs': 10})
    hm_rows.append(row)

# Column widths
cw = [Inches(1.2)] + [Inches(0.7)] * 14 + [Inches(0.6)]
colored_cell_table(sl, Inches(0.3), Inches(1.2), hm_rows, cw)

# Legend
y = Inches(5.6)
add_textbox(sl, Inches(0.8), y, Inches(1.0), Inches(0.3), "Legend:", font_size=10, bold=True, color=DARK)
for level in range(4):
    x = Inches(1.8 + level * 2.0)
    badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.6), Inches(0.3))
    badge.fill.solid()
    badge.fill.fore_color.rgb = SCORE_COLORS[level]
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"{level} = {SCORE_LABELS[level]}"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Observations
add_textbox(sl, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
            "Observations: ECCO and ISSM lead in overall capability (process representation, verification & validation, mission support).\n"
            "Common gaps: machine learning / artificial intelligence integration, stakeholder adoption, and computational performance across most models.",
            font_size=12, color=GRAY500)

slide_number_footer(sl, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MODEL DEEP DIVES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), INDIGO600)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Model Deep Dives: ECCO & ISSM", font_size=30, bold=True, color=DARK)

for col, (model, mc, strengths, gaps) in enumerate([
    ("ECCO", PURPLE500,
     ["Process Representation (3): State-of-the-art ocean general circulation model",
      "Observational Constraint (3): 14+ observation types assimilated",
      "Retrospective Analysis (3): version 4 release 4 reanalysis 1992-2017",
      "Verification & Validation Framework (3): Extensive intercomparison participation",
      "Mission Support (3): Observing System Simulation Experiments for GRACE Follow-On, SWOT, Sentinel-6"],
     ["Machine Learning / Artificial Intelligence Integration (1): Research-stage only",
      "Stakeholder Adoption (2): Growing but not yet operational"]),
    ("ISSM", CYAN500,
     ["Process Representation (3): Full ice-sheet dynamics + solid Earth",
      "Spatial Resolution (3): Adaptive mesh, sub-kilometer at margins",
      "Process Coupling (3): 3-way ice-solid Earth-sea level",
      "Verification & Validation Framework (3): Ice Sheet Model Intercomparison Project for Coupled Model Intercomparison Project Phase 6, initial state Model Intercomparison Project benchmarking",
      "Interoperability (3): Open source, active community"],
     ["Machine Learning / Artificial Intelligence Integration (1): Exploratory only",
      "Retrospective Analysis (2): Not continuous reanalysis",
      "Stakeholder Adoption (2): Expert interpretation required"]),
]):
    x = Inches(0.5 + col * 6.3)

    # Model header
    add_accent_bar(sl, x, Inches(1.1), Inches(5.8), Inches(0.05), mc)
    add_textbox(sl, x, Inches(1.3), Inches(5.8), Inches(0.5),
                model, font_size=22, bold=True, color=mc)

    # Compute avg
    t1 = mcl_data.get(model, {}).get('tier1', {})
    scores = [v.get('score', 0) for v in t1.values()]
    avg = sum(scores) / len(scores) if scores else 0
    add_textbox(sl, x + Inches(3.5), Inches(1.3), Inches(2.3), Inches(0.5),
                f"Tier 1 Average: {avg:.1f}/3.0", font_size=14, bold=True, color=DARK, alignment=PP_ALIGN.RIGHT)

    # Strengths
    add_textbox(sl, x, Inches(1.9), Inches(5.8), Inches(0.3),
                "Strengths (Level 3)", font_size=13, bold=True, color=GREEN600)
    add_bullet_list(sl, x, Inches(2.3), Inches(5.8), Inches(2.0),
                    strengths, font_size=11, color=GRAY700)

    # Gaps
    gap_y = Inches(4.4) if len(strengths) <= 5 else Inches(4.8)
    add_textbox(sl, x, gap_y, Inches(5.8), Inches(0.3),
                "Improvement Areas", font_size=13, bold=True, color=AMBER500)
    add_bullet_list(sl, x, gap_y + Inches(0.35), Inches(5.8), Inches(1.5),
                    gaps, font_size=11, color=GRAY700)

slide_number_footer(sl, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GAP ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), RED500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Gap Analysis: Common Development Needs", font_size=30, bold=True, color=DARK)

# Compute gap counts per dimension
gap_data = []
for dim_id, dim_label in TIER1_DIMS:
    count_low = sum(1 for m in MODELS if mcl_data.get(m, {}).get('tier1', {}).get(dim_id, {}).get('score', 0) < 2)
    avg = sum(mcl_data.get(m, {}).get('tier1', {}).get(dim_id, {}).get('score', 0) for m in MODELS) / len(MODELS)
    gap_data.append((dim_label, count_low, avg))

gap_data.sort(key=lambda x: -x[1])

# Table
gh = [
    {'text': 'Dimension', 'bg': RGBColor(0x7F, 0x1D, 0x1D), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Models < Level 2', 'bg': RGBColor(0x7F, 0x1D, 0x1D), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Mean Score', 'bg': RGBColor(0x7F, 0x1D, 0x1D), 'fg': WHITE, 'bold': True, 'fs': 11},
    {'text': 'Priority', 'bg': RGBColor(0x7F, 0x1D, 0x1D), 'fg': WHITE, 'bold': True, 'fs': 11},
]
gap_rows = [gh]
for label, count, avg in gap_data[:10]:
    priority = 'HIGH' if count >= 5 else ('MEDIUM' if count >= 3 else 'LOW')
    pc = RED500 if priority == 'HIGH' else (AMBER500 if priority == 'MEDIUM' else GREEN600)
    gap_rows.append([
        {'text': label, 'fg': DARK, 'bold': False, 'fs': 11, 'align': PP_ALIGN.LEFT},
        {'text': f'{count}/{len(MODELS)}', 'fg': RED500 if count >= 5 else GRAY700, 'bold': False, 'fs': 12},
        {'text': f'{avg:.1f}', 'fg': GRAY700, 'bold': False, 'fs': 12},
        {'text': priority, 'fg': pc, 'bold': False, 'fs': 11},
    ])

colored_cell_table(sl, Inches(0.5), Inches(1.3), gap_rows, [Inches(3.0), Inches(2.0), Inches(1.5), Inches(1.5)])

# Recommendations
add_textbox(sl, Inches(8.5), Inches(1.3), Inches(4.5), Inches(0.5),
            "Strategic Recommendations", font_size=16, bold=True, color=RED500)
recs = [
    "Machine Learning / Artificial Intelligence Integration: Most models at Level 0-1. Cross-model machine learning investment would benefit all teams.",
    "Stakeholder Adoption: Bridge from research to operational use. Needs dedicated applications engagement.",
    "Computational Performance: Graphics Processing Unit / cloud portability is a common barrier. Shared infrastructure investment.",
    "Retrospective Analysis: Several models lack long-term reanalysis. Opportunity for coupled reanalysis products."
]
add_bullet_list(sl, Inches(8.5), Inches(1.9), Inches(4.5), Inches(4.5),
                recs, font_size=11, color=GRAY700, bold_prefix=True)

slide_number_footer(sl, 13, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DASHBOARD ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, WHITE)
add_accent_bar(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.06), GRAY500)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Dashboard Architecture", font_size=30, bold=True, color=DARK)

# Tech stack
add_textbox(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
            "Technology Stack", font_size=18, bold=True, color=GRAY700)
tech_items = [
    "React (Create React App) with React Router for single-page application navigation",
    "Tailwind Cascading Style Sheets for responsive, utility-first styling",
    "Recharts for interactive data visualizations",
    "D3.js + TopoJSON for force-directed graphs and maps",
    "JavaScript Object Notation data files with dynamic imports per model",
    "Lucide React for consistent iconography",
]
add_bullet_list(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0),
                tech_items, font_size=12, color=GRAY700)

# Key features
add_textbox(sl, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
            "Key Features", font_size=18, bold=True, color=GRAY700)
features = [
    "Context-aware navigation: JPL's Earth Modeling Enterprise ↔ JPL Earth Observation Enterprise with dynamic branding",
    "Generic + model-specific page architecture",
    "Cross-model comparison: network graph, heatmap, radar chart",
    "Per-model drill-down: citations, geography, domains, uncertainty quantification, maturity",
    "Model Capability Level dashboard: filterable radar, heatmap, gap analysis",
    "Evidence-linked scoring: click any cell for justification",
]
add_bullet_list(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0),
                features, font_size=12, color=GRAY700)

# Route structure
add_textbox(sl, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
            "Route Structure", font_size=16, bold=True, color=GRAY700)
routes = [
    "/science-model-dashboard — JEME main dashboard (8 models)",
    "/science-model-dashboard/JEOE — JEOE main dashboard (missions)",
    "/science-model-dashboard/{model} — Model-specific dashboard",
    "/science-model-dashboard/{model}/citations | geographic-impact | research-domains | uncertainty | maturity",
    "/science-model-dashboard/earth-system — Earth System Interconnections",
    "/science-model-dashboard/model-maturity — Cross-model MCL comparison",
]
add_bullet_list(sl, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
                routes, font_size=10, color=GRAY500)

slide_number_footer(sl, 14, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PATH FORWARD & SUMMARY
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(sl, RGBColor(0x0F, 0x17, 0x2A))

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Path Forward & Summary", font_size=32, bold=True, color=WHITE)

# Path forward
add_accent_bar(sl, Inches(0.8), Inches(1.3), Inches(5.0), Inches(0.04), BLUE500)
add_textbox(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Next Steps", font_size=20, bold=True, color=BLUE500)
next_items = [
    "Large Language Model-enhanced Model Capability Level scoring: auto-classify papers against dimensions, generate evidence summaries",
    "Periodic reassessment: establish annual Model Capability Level review cycle with model teams",
    "Community adoption: engage external evaluators for independent scoring",
    "Expand to coupled assessments: evaluate inter-model coupling maturity",
    "Extend JPL Earth Observation Enterprise: add more NASA missions as data becomes available",
    "Historical tracking: Model Capability Level score timelines to show model evolution"
]
add_bullet_list(sl, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5),
                next_items, font_size=12, color=RGBColor(0xBF, 0xDB, 0xFE))

# Key takeaways
add_accent_bar(sl, Inches(7.0), Inches(1.3), Inches(5.0), Inches(0.04), GREEN500)
add_textbox(sl, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
            "Key Takeaways", font_size=20, bold=True, color=GREEN500)
takeaways = [
    "Unified dashboard tracks citation impact across 8 models and 2 missions",
    "Multi-agent pipeline ensures data quality with calibrated uncertainty",
    "Earth system classification reveals cross-model research connections",
    "Model Capability Level framework provides systematic, evidence-based capability assessment",
    "Common gaps identified: machine learning / artificial intelligence integration and stakeholder adoption",
    "Framework is living and extensible — designed for periodic reassessment"
]
add_bullet_list(sl, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.5),
                takeaways, font_size=12, color=RGBColor(0xBB, 0xF7, 0xD0))

# Bottom
add_accent_bar(sl, Inches(0), Inches(6.2), Inches(13.33), Inches(0.03), RGBColor(0x33, 0x44, 0x66))
add_textbox(sl, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
            "Dashboard:  http://34.31.165.25:3000/science-model-dashboard/",
            font_size=12, color=GRAY400)
add_textbox(sl, Inches(0.8), Inches(6.75), Inches(11), Inches(0.5),
            "Thank you — Questions?", font_size=18, bold=True, color=WHITE)

slide_number_footer(sl, 15, TOTAL_SLIDES)


# ── Save ───────────────────────────────────────────────────────────────────
output_path = os.path.join(project_dir, 'JEME_JEOE_Technical_Seminar.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
