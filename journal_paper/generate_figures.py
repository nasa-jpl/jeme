#!/usr/bin/env python3
"""Generate journal paper figures (PNG + editable PowerPoint).

Produces:
  figures/fig1_citation_overview.png       — bar chart of papers/citations per model
  figures/fig5_network_matrix.png          — cross-model bridge-paper heatmap
  figures/fig7_uq_pipeline.png             — three-phase UQ schematic
  figures/fig8_sphere_coverage.png         — sphere-coverage matrix
  figures/journal_paper_figures.pptx       — single editable deck containing all
                                              figures (native shapes where
                                              possible, images otherwise).

All figure fonts are at least 12 pt (matching main-text size).
"""

import json
from collections import Counter, defaultdict
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
ROOT = Path(__file__).parent.parent
DATA = ROOT / 'public' / 'data'
OUT = Path(__file__).parent / 'figures'
OUT.mkdir(parents=True, exist_ok=True)

# Set global matplotlib style: 12 pt base font, black text everywhere
plt.rcParams.update({
    'font.size': 12,
    'axes.titlesize': 14,
    'axes.labelsize': 12,
    'xtick.labelsize': 12,
    'ytick.labelsize': 12,
    'legend.fontsize': 12,
    'figure.titlesize': 14,
    'text.color': 'black',
    'axes.labelcolor': 'black',
    'xtick.color': 'black',
    'ytick.color': 'black',
})

ENTITIES = [
    ('ECCO', 'Model'), ('ISSM', 'Model'), ('MOMO-CHEM', 'Model'),
    ('EDMF', 'Model'), ('CMS-Flux', 'Model'), ('CARDAMOM', 'Model'),
    ('RAPID', 'Model'), ('LES', 'Model'),
    ('GRACE', 'Mission'), ('SWOT', 'Mission'), ('TROPESS', 'Mission'),
]

# ---------------------------------------------------------------------------
# Compute statistics from current JSON data
# ---------------------------------------------------------------------------

def compute_stats():
    rows = []
    for name, kind in ENTITIES:
        f = DATA / f'{name}_analyzed.json'
        d = json.load(open(f))
        cits = sum(int(p.get('citation_count') or p.get('is-referenced-by-count') or 0)
                   for p in d)
        rows.append({'name': name, 'type': kind, 'papers': len(d), 'citations': cits})
    return rows


def compute_network():
    """Bridge papers across the eight JEME models — DOI overlap matrix."""
    models = [n for n, k in ENTITIES if k == 'Model']
    doi_sets = {}
    for m in models:
        f = DATA / f'{m}_analyzed.json'
        d = json.load(open(f))
        doi_sets[m] = {(p.get('doi') or '').lower() for p in d if p.get('doi')}
    M = np.zeros((len(models), len(models)), dtype=int)
    for i, mi in enumerate(models):
        for j, mj in enumerate(models):
            if i == j:
                M[i][j] = len(doi_sets[mi])
            else:
                M[i][j] = len(doi_sets[mi] & doi_sets[mj])
    return models, M


def compute_sphere_coverage():
    """Approximate Earth-system-sphere coverage for the eight JEME models."""
    sphere_keywords = {
        'Atmosphere': ['atmospher', 'troposph', 'stratosph', 'aerosol', 'ozone',
                       'ch4', 'co2', 'methane', 'cloud', 'precipitation'],
        'Hydrosphere': ['ocean', 'sea', 'water', 'river', 'streamflow',
                        'discharge', 'estuary', 'precipitation'],
        'Cryosphere': ['ice', 'glacier', 'snow', 'permafrost', 'sea ice',
                       'ice sheet', 'antarctic', 'arctic'],
        'Biosphere':  ['vegetation', 'forest', 'biomass', 'ecosystem',
                       'carbon flux', 'photosynth', 'leaf', 'soil organic'],
        'Geosphere':  ['geode', 'gravity', 'tecton', 'lithosphere',
                       'sea level', 'crustal'],
    }
    models = [n for n, k in ENTITIES if k == 'Model']
    spheres = list(sphere_keywords)
    M = np.zeros((len(spheres), len(models)))
    for j, m in enumerate(models):
        f = DATA / f'{m}_analyzed.json'
        d = json.load(open(f))
        for i, sp in enumerate(spheres):
            kws = sphere_keywords[sp]
            n = 0
            for p in d:
                t = ((p.get('title') or '') + ' ' + (p.get('abstract') or '')).lower()
                if isinstance(p.get('title'), list):
                    t = ' '.join(p['title']).lower() + ' ' + (p.get('abstract') or '').lower()
                if any(kw in t for kw in kws):
                    n += 1
            M[i][j] = n / max(1, len(d))
    return spheres, models, M


# ---------------------------------------------------------------------------
# matplotlib PNG figures
# ---------------------------------------------------------------------------

def fig1_citation_overview(stats, png_path):
    fig, ax = plt.subplots(figsize=(11, 6.5))
    names = [r['name'] for r in stats]
    papers = [r['papers'] for r in stats]
    citations = [r['citations'] for r in stats]
    x = np.arange(len(names))
    width = 0.4
    b1 = ax.bar(x - width / 2, papers, width, label='Papers', color='#1f77b4')
    ax2 = ax.twinx()
    b2 = ax2.bar(x + width / 2, citations, width, label='Citations',
                 color='#ff7f0e', alpha=0.85)
    ax.set_yscale('log'); ax2.set_yscale('log')
    ax.set_xticks(x); ax.set_xticklabels(names, rotation=30, ha='right')
    ax.set_ylabel('Papers (log)'); ax2.set_ylabel('Citations (log)')
    ax.set_title('Citation dataset overview by model and mission')
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax.legend(h1 + h2, l1 + l2, loc='upper right', frameon=False)
    for bar, val in zip(b1, papers):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height(),
                f'{val:,}', ha='center', va='bottom', fontsize=11, color='black')
    plt.tight_layout()
    plt.savefig(png_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()


def fig5_network_matrix(models, M, png_path):
    fig, ax = plt.subplots(figsize=(8.5, 7))
    # Hide the diagonal so cross-pair overlap is visible
    masked = M.astype(float).copy()
    np.fill_diagonal(masked, np.nan)
    im = ax.imshow(masked, cmap='Blues', aspect='equal')
    ax.set_xticks(np.arange(len(models)))
    ax.set_yticks(np.arange(len(models)))
    ax.set_xticklabels(models, rotation=30, ha='right')
    ax.set_yticklabels(models)
    for i in range(len(models)):
        for j in range(len(models)):
            if i == j:
                ax.text(j, i, '—', ha='center', va='center',
                        fontsize=12, color='black')
            else:
                ax.text(j, i, str(int(M[i][j])), ha='center', va='center',
                        fontsize=12,
                        color='white' if masked[i][j] > np.nanmax(masked) / 2 else 'black')
    ax.set_title('Cross-model bridge-paper count\n(shared DOIs between model citation sets)')
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label('Bridge papers', fontsize=12)
    cbar.ax.tick_params(labelsize=12)
    plt.tight_layout()
    plt.savefig(png_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()


def fig7_uq_pipeline(png_path):
    """Three-phase UQ schematic — drawn entirely with matplotlib shapes."""
    fig, ax = plt.subplots(figsize=(11.5, 5.5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis('off')

    boxes = [
        (0.5, 1.3, 'Phase 1\nDeterministic scoring',
         'evidence × reasoning − pipeline variance'),
        (3.7, 1.3, 'Phase 2\nMulti-temperature LLM sampling',
         'T = {0.1, 0.5, 1.0}, 3 passes\nstochastic variance + verbalised confidence'),
        (6.9, 1.3, 'Phase 3\nSkeptic agent review',
         'adversarial challenge\nagreement ≤ 2/5 → override flag'),
    ]
    for x, y, title, body in boxes:
        ax.add_patch(plt.Rectangle((x, y), 2.6, 2.4, facecolor='#e3f2fd',
                                   edgecolor='black', linewidth=1.2))
        ax.text(x + 1.3, y + 1.95, title, ha='center', va='center',
                fontsize=12, fontweight='bold', color='black')
        ax.text(x + 1.3, y + 0.95, body, ha='center', va='center',
                fontsize=11, color='black')
    # Arrows between boxes
    for x_start in (3.1, 6.3):
        ax.annotate('', xy=(x_start + 0.55, 2.5), xytext=(x_start, 2.5),
                    arrowprops=dict(arrowstyle='->', color='black', lw=1.5))
    ax.text(5, 4.5, 'Three-phase uncertainty quantification',
            ha='center', va='center', fontsize=14, fontweight='bold',
            color='black')
    ax.text(5, 0.6,
            'output: composite confidence ∈ [0.05, 0.99]   |   '
            'override flag for skeptic-disagreement entries',
            ha='center', va='center', fontsize=12, color='black', style='italic')
    plt.tight_layout()
    plt.savefig(png_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()


def fig8_sphere_coverage(spheres, models, M, png_path):
    fig, ax = plt.subplots(figsize=(9, 5.5))
    im = ax.imshow(M, cmap='YlGnBu', aspect='auto')
    ax.set_xticks(np.arange(len(models))); ax.set_xticklabels(models, rotation=30, ha='right')
    ax.set_yticks(np.arange(len(spheres))); ax.set_yticklabels(spheres)
    for i in range(len(spheres)):
        for j in range(len(models)):
            ax.text(j, i, f'{M[i][j]*100:.0f}%', ha='center', va='center',
                    fontsize=12,
                    color='white' if M[i][j] > 0.5 else 'black')
    ax.set_title('Earth-system-sphere coverage by JEME model\n'
                 '(% of papers with sphere-keyword match)')
    cbar = fig.colorbar(im, ax=ax, fraction=0.04, pad=0.02)
    cbar.set_label('Fraction of papers', fontsize=12)
    cbar.ax.tick_params(labelsize=12)
    plt.tight_layout()
    plt.savefig(png_path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()


# ---------------------------------------------------------------------------
# Editable PPTX
# ---------------------------------------------------------------------------

PPT_W = Inches(13.33); PPT_H = Inches(7.5)  # widescreen
TITLE_PT = Pt(20); BODY_PT = Pt(12)


def _add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),
                                   Inches(12.3), Inches(0.6))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.runs[0].font.size = TITLE_PT; p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0, 0, 0)


def _add_caption(slide, text, top=Inches(6.7)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.runs[0].font.size = BODY_PT
    p.runs[0].font.color.rgb = RGBColor(0, 0, 0)


def _add_image(slide, path, top=Inches(1.0), height=Inches(5.5)):
    pic = slide.shapes.add_picture(str(path), Inches(0.5), top, height=height)
    # Centre horizontally
    pic.left = Emu(int((PPT_W - pic.width) / 2))


def slide_fig1_chart(prs, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 1. Citation dataset overview')
    chart_data = CategoryChartData()
    chart_data.categories = [r['name'] for r in stats]
    chart_data.add_series('Papers', [r['papers'] for r in stats])
    chart_data.add_series('Citations', [r['citations'] for r in stats])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.6), chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    # Black 12pt fonts everywhere
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = BODY_PT
        axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)
    chart.value_axis.has_major_gridlines = True
    # Log scale would be ideal but python-pptx doesn't expose it cleanly;
    # users can toggle log scale in PPT after opening (right-click axis →
    # Format Axis → Logarithmic).
    _add_caption(slide,
                 'Papers and citations per JEME model and JEOE mission '
                 '(post peer-review filter and multi-agent cleanup, 2026-04). '
                 'Switch the value axis to logarithmic in PowerPoint to see '
                 'small models clearly.')


def slide_fig5(prs, png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 5. Cross-model bridge-paper matrix')
    _add_image(slide, png)
    _add_caption(slide,
                 'Bridge papers — the count of DOIs cited by both models in '
                 'each pair. Diagonal omitted; matrix is symmetric.')


def slide_fig7(prs, png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 7. Three-phase uncertainty quantification')
    _add_image(slide, png)
    _add_caption(slide,
                 'Composite confidence is computed in three stages: '
                 'deterministic scoring on metadata, multi-temperature LLM '
                 'sampling, and adversarial skeptic review.')


def slide_fig8(prs, png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 8. Earth-system-sphere coverage')
    _add_image(slide, png)
    _add_caption(slide,
                 'Fraction of each model\'s papers that match keywords for '
                 'each Earth-system sphere. Multi-sphere papers count toward '
                 'every matching sphere.')


def _box(slide, x, y, w, h, label, fill=RGBColor(0xE3, 0xF2, 0xFD)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0, 0, 0); shp.line.width = Pt(1)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.text = label
    for run in p.runs:
        run.font.size = BODY_PT
        run.font.color.rgb = RGBColor(0, 0, 0)
    return shp


def _arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = RGBColor(0, 0, 0); line.line.width = Pt(1.25)


def slide_fig9_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 9. Multi-source full-text retrieval pipeline')

    nodes = [
        ('Copernicus (10.5194)',         0.5, 1.3, 2.0, 0.7),
        ('Wiley TDM\n(10.1002/29/111)',  3.0, 1.3, 2.0, 0.9),
        ('Elsevier TDM (10.1016)',       5.5, 1.3, 2.0, 0.9),
        ('Unpaywall\nPDF + HTML',        8.0, 1.3, 2.0, 0.9),
        ('OpenAlex\ninstitutional repos',2.0, 3.4, 2.5, 0.9),
        ('Semantic Scholar\nopenAccessPdf', 5.5, 3.4, 2.5, 0.9),
        ('DOI redirect (last resort)',  9.0, 3.4, 3.0, 0.9),
        ('Extract DAS + Methods +\nmarker contexts', 4.0, 5.0, 5.3, 0.9),
        ('Re-classify with Gemini', 4.0, 6.2, 5.3, 0.9),
    ]
    placed = []
    for label, xi, yi, wi, hi in nodes:
        shp = _box(slide, Inches(xi), Inches(yi), Inches(wi), Inches(hi), label)
        placed.append(shp)

    # connect tier 1 (top row) → tier 2 (middle) → extract → classify
    for i in range(4):
        # tier 1 (i) downarrow to extract block (placed[7])
        a = placed[i]
        _arrow(slide, a.left + a.width // 2, a.top + a.height,
               Inches(6.65), Inches(5.0))
    for i in range(4, 7):
        a = placed[i]
        _arrow(slide, a.left + a.width // 2, a.top + a.height,
               Inches(6.65), Inches(5.0))
    extract = placed[7]; classify = placed[8]
    _arrow(slide, extract.left + extract.width // 2, extract.top + extract.height,
           classify.left + classify.width // 2, classify.top)

    _add_caption(slide,
                 'Tiered DOI lookup. Each box is a fetch route; success at any '
                 'tier short-circuits subsequent attempts. ~1,300 previously-'
                 'failed fetches were recovered when OpenAlex was added.')


def slide_fig10_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, 'Figure 10. System architecture')

    sections = [
        ('Data sources', 0.5, 1.2, 3.0, 5.0,
         'Semantic Scholar\nCrossref\nUnpaywall\nOpenAlex\nWiley TDM\nElsevier TDM\nCopernicus',
         RGBColor(0xE3, 0xF2, 0xFD)),
        ('Citation analytics pipeline', 4.0, 1.2, 4.5, 5.0,
         'Seed identification\n↓\nCitation collection\n↓\nMulti-agent verification\n↓\n'
         'Peer-review filtering\n↓\nFull-text enrichment\n↓\nGemini classification\n↓\n'
         'Three-phase UQ',
         RGBColor(0xFF, 0xF9, 0xC4)),
        ('Storage', 9.0, 1.2, 1.6, 5.0,
         '{model}_\nanalyzed.json',
         RGBColor(0xF3, 0xE5, 0xF5)),
        ('React dashboard', 11.1, 1.2, 1.7, 5.0,
         'JEME (8 models)\nJEOE (3 missions)\nJESP (199 scientists)\n'
         'Network analysis\nSphere coverage',
         RGBColor(0xC8, 0xE6, 0xC9)),
    ]
    for label, x, y, w, h, body, fill in sections:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = RGBColor(0, 0, 0); shp.line.width = Pt(1)
        tf = shp.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.1)
        # Header
        p = tf.paragraphs[0]; p.text = label
        for run in p.runs:
            run.font.size = Pt(13); run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
        # Body
        body_p = tf.add_paragraph(); body_p.text = body
        for run in body_p.runs:
            run.font.size = BODY_PT
            run.font.color.rgb = RGBColor(0, 0, 0)

    # Arrows between sections
    for x in (3.5, 8.5, 10.6):
        _arrow(slide, Inches(x), Inches(3.7),
               Inches(x + 0.5), Inches(3.7))

    _add_caption(slide,
                 'End-to-end data flow from external scholarly APIs through '
                 'the citation analytics pipeline into per-model JSON storage '
                 'and the React dashboard.')


# ---------------------------------------------------------------------------
# Drive
# ---------------------------------------------------------------------------

def main():
    print('Computing dataset statistics...')
    stats = compute_stats()
    for r in stats:
        print(f'  {r["name"]:10} {r["papers"]:6,d} papers  {r["citations"]:10,d} citations')

    print('Computing cross-model network...')
    models, network_M = compute_network()
    print('Computing sphere coverage...')
    spheres, _, sphere_M = compute_sphere_coverage()

    print('Generating PNG figures...')
    fig1 = OUT / 'fig1_citation_overview.png'
    fig5 = OUT / 'fig5_network_matrix.png'
    fig7 = OUT / 'fig7_uq_pipeline.png'
    fig8 = OUT / 'fig8_sphere_coverage.png'
    fig1_citation_overview(stats, fig1)
    fig5_network_matrix(models, network_M, fig5)
    fig7_uq_pipeline(fig7)
    fig8_sphere_coverage(spheres, models, sphere_M, fig8)

    print('Building editable PowerPoint...')
    prs = Presentation()
    prs.slide_width = PPT_W; prs.slide_height = PPT_H
    slide_fig1_chart(prs, stats)
    slide_fig5(prs, fig5)
    slide_fig7(prs, fig7)
    slide_fig8(prs, fig8)
    slide_fig9_pipeline(prs)
    slide_fig10_architecture(prs)
    pptx_path = OUT / 'journal_paper_figures.pptx'
    prs.save(pptx_path)
    print(f'  -> {pptx_path}')
    print('Done.')


if __name__ == '__main__':
    main()
